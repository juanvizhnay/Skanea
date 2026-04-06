from __future__ import annotations

import io
import os
import csv
import chardet
from typing import Any, Dict

import fitz  # PyMuPDF
import pandas as pd
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from fastapi import FastAPI, File, HTTPException, UploadFile
import sys
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from dotenv import load_dotenv

from .cache import compute_sha256, load_cache, write_cache
from .pdf_extract import detect_native_pdf, extract_pdf_native, extract_pdf_scanned
from .ocr import ocr_image_to_text


MAX_SIZE_BYTES = int(os.getenv("MAX_UPLOAD_MB", "100")) * 1024 * 1024
ALLOWED_PDF_TYPES = {"application/pdf"}
ALLOWED_IMAGE_TYPES = {"image/png", "image/jpeg", "image/jpg", "image/webp"}
ALLOWED_DOCX_TYPES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
}
ALLOWED_PPTX_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
}
ALLOWED_XLSX_TYPES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
}
ALLOWED_CSV_TYPES = {"text/csv", "application/csv"}
ALLOWED_TXT_TYPES = {"text/plain", "text/txt"}


load_dotenv()
try:
    if os.getenv("SK_EXTRACT_DEBUG", "0") == "1":
        print(f"[INIT] Python exec={sys.executable} version={sys.version}")
except Exception:
    pass
app = FastAPI(title="Skanea Extract Service", version="0.1.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "DELETE", "OPTIONS"],
    allow_headers=["*"],
    max_age=600,
)


@app.get("/health")
def health() -> Dict[str, Any]:
    return {"ok": True}


@app.get("/debug/ocr")
def debug_ocr() -> Dict[str, Any]:
    """Debug endpoint to check OCR status"""
    try:
        from .ocr import diagnose_ocr_backends

        diagnosis = diagnose_ocr_backends()
        return {"ok": True, "backends": diagnosis}
    except Exception as e:
        return {"ok": False, "error": str(e)}


# Warmup: Precargar modelos pesados al inicio (pix2text para Math OCR)
@app.on_event("startup")
def _warmup_math_ocr_models() -> None:
    """
    Precarga pix2text al iniciar el servidor para que la primera imagen
    matematica no demore 30 segundos. Toma ~26s al inicio pero despues
    todas las imagenes matematicas responden en ~3-5s.
    """
    import time
    t_start = time.time()
    print("[WARMUP] Precargando modelos de Math OCR (pix2text)...")
    print("[WARMUP] Esto tomara ~25-30 segundos pero solo pasa una vez...")
    
    try:
        from .math_ocr import get_pix2text
        
        # Inicializar pix2text (carga TensorFlow, PyTorch, y 5+ modelos ML)
        p2t = get_pix2text()
        
        t_end = time.time()
        elapsed = t_end - t_start
        print(f"[WARMUP] OK - pix2text precargado exitosamente en {elapsed:.1f}s")
        print("[WARMUP] El servidor esta listo - Math OCR respondera rapido ahora!")
        
    except Exception as e:
        t_end = time.time()
        elapsed = t_end - t_start
        print(f"[WARMUP] ERROR - Fallo precargando pix2text (tomo {elapsed:.1f}s): {e}")
        print("[WARMUP] Math OCR aun funcionara pero sera lento la primera vez")


def read_upload(file: UploadFile) -> bytes:
    content = file.file.read()
    if not content:
        raise HTTPException(
            status_code=400,
            detail={"ok": False, "error": "Archivo vacío", "code": "E_EMPTY"},
        )
    if len(content) > MAX_SIZE_BYTES:
        raise HTTPException(
            status_code=413,
            detail={
                "ok": False,
                "error": "Archivo demasiado grande",
                "code": "E_FILE_TOO_LARGE",
            },
        )
    return content


def error_response(status_code: int, code: str, message: str) -> JSONResponse:
    return JSONResponse(
        status_code=status_code, content={"ok": False, "error": message, "code": code}
    )


@app.post("/extract/pdf")
def extract_pdf(file: UploadFile = File(...)) -> Dict[str, Any]:
    if file.content_type not in ALLOWED_PDF_TYPES:
        return error_response(
            415, "E_UNSUPPORTED_TYPE", f"Tipo no soportado: {file.content_type}"
        )
    content = read_upload(file)
    sha = compute_sha256(content)
    cached = load_cache(sha)
    if cached:
        cached.setdefault("meta", {})["cache_hit"] = True
        return cached
    try:
        doc = fitz.open(stream=content, filetype="pdf")
    except Exception:
        return error_response(400, "E_PDF_OPEN", "No se pudo abrir el PDF")
    try:
        is_native = detect_native_pdf(doc)
        if is_native:
            pages, avg_conf = extract_pdf_native(doc)
        else:
            pages, avg_conf = extract_pdf_scanned(doc, dpi=300)
    finally:
        doc.close()
    full_text = "\n\n".join(p.get("text", "") for p in pages)
    resp = {
        "ok": True,
        "is_native": is_native,
        "hash": sha,
        "meta": {"pages": len(pages), "lang": "auto", "cache_hit": False},
        "confidence": avg_conf,
        "pages": pages,
        "full_text": full_text,
        "warnings": [],
    }
    write_cache(sha, resp)
    return resp


@app.post("/extract/image")
def extract_image(file: UploadFile = File(...)) -> Dict[str, Any]:
    if file.content_type not in ALLOWED_IMAGE_TYPES:
        return error_response(
            415, "E_UNSUPPORTED_TYPE", f"Tipo no soportado: {file.content_type}"
        )
    content = read_upload(file)
    sha = compute_sha256(content)
    cached = load_cache(sha)
    if cached:
        cached.setdefault("meta", {})["cache_hit"] = True
        return cached
    try:
        from PIL import Image

        img = Image.open(io.BytesIO(content))
        res = ocr_image_to_text(img)
        pages = [
            {
                "number": 1,
                "text": res.get("text", ""),
                "confidence": float(res.get("confidence", 0.0)),
            }
        ]
        avg_conf = float(res.get("confidence", 0.0))
        full_text = res.get("text", "")
        resp = {
            "ok": True,
            "hash": sha,
            "meta": {"pages": 1, "lang": "auto", "cache_hit": False},
            "confidence": avg_conf,
            "pages": pages,
            "full_text": full_text,
            "warnings": [],
        }
        write_cache(sha, resp)
        return resp
    except Exception:
        return error_response(400, "E_IMAGE_OPEN", "No se pudo procesar la imagen")


@app.post("/extract/docx")
def extract_docx(file: UploadFile = File(...)) -> Dict[str, Any]:
    if file.content_type not in ALLOWED_DOCX_TYPES:
        return error_response(
            415, "E_UNSUPPORTED_TYPE", f"Tipo no soportado: {file.content_type}"
        )
    content = read_upload(file)
    sha = compute_sha256(content)
    cached = load_cache(sha)
    if cached:
        cached.setdefault("meta", {})["cache_hit"] = True
        return cached
    try:
        doc = Document(io.BytesIO(content))
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())

        full_text = "\n".join(paragraphs)
        pages = [{"number": 1, "text": full_text, "confidence": 1.0}]

        resp = {
            "ok": True,
            "hash": sha,
            "meta": {
                "pages": 1,
                "paragraphs": len(paragraphs),
                "lang": "auto",
                "cache_hit": False,
            },
            "confidence": 1.0,
            "pages": pages,
            "full_text": full_text,
            "warnings": [],
        }
        write_cache(sha, resp)
        return resp
    except Exception:
        return error_response(
            400, "E_DOCX_OPEN", "No se pudo procesar el documento DOCX"
        )


@app.post("/extract/xlsx")
def extract_xlsx(file: UploadFile = File(...)) -> Dict[str, Any]:
    if file.content_type not in ALLOWED_XLSX_TYPES:
        return error_response(
            415, "E_UNSUPPORTED_TYPE", f"Tipo no soportado: {file.content_type}"
        )
    content = read_upload(file)
    sha = compute_sha256(content)
    cached = load_cache(sha)
    if cached:
        cached.setdefault("meta", {})["cache_hit"] = True
        return cached
    try:
        workbook = load_workbook(io.BytesIO(content), read_only=True)
        all_sheets_text = []
        pages = []

        for i, sheet_name in enumerate(workbook.sheetnames):
            sheet = workbook[sheet_name]
            sheet_text = []

            for row in sheet.iter_rows(values_only=True):
                row_text = []
                for cell in row:
                    if cell is not None:
                        row_text.append(str(cell))
                if row_text:
                    sheet_text.append("\t".join(row_text))

            sheet_content = "\n".join(sheet_text)
            if sheet_content.strip():
                pages.append(
                    {
                        "number": i + 1,
                        "sheet_name": sheet_name,
                        "text": sheet_content,
                        "confidence": 1.0,
                    }
                )
                all_sheets_text.append(f"== Hoja: {sheet_name} ==\n{sheet_content}")

        workbook.close()
        full_text = "\n\n".join(all_sheets_text)

        resp = {
            "ok": True,
            "hash": sha,
            "meta": {
                "pages": len(pages),
                "sheets": len(workbook.sheetnames),
                "lang": "auto",
                "cache_hit": False,
            },
            "confidence": 1.0,
            "pages": pages,
            "full_text": full_text,
            "warnings": [],
        }
        write_cache(sha, resp)
        return resp
    except Exception:
        return error_response(
            400, "E_XLSX_OPEN", "No se pudo procesar el archivo Excel"
        )


@app.post("/extract/csv")
def extract_csv(file: UploadFile = File(...)) -> Dict[str, Any]:
    if file.content_type not in ALLOWED_CSV_TYPES:
        return error_response(
            415, "E_UNSUPPORTED_TYPE", f"Tipo no soportado: {file.content_type}"
        )
    content = read_upload(file)
    sha = compute_sha256(content)
    cached = load_cache(sha)
    if cached:
        cached.setdefault("meta", {})["cache_hit"] = True
        return cached
    try:
        # Detectar encoding
        detected = chardet.detect(content)
        encoding = detected.get("encoding", "utf-8")

        text_content = content.decode(encoding)
        csv_reader = csv.reader(io.StringIO(text_content))

        rows = []
        headers = None
        for i, row in enumerate(csv_reader):
            if i == 0:
                headers = row
                rows.append("\t".join(row))  # Primera fila como headers
            else:
                rows.append("\t".join(row))

        full_text = "\n".join(rows)
        pages = [{"number": 1, "text": full_text, "confidence": 1.0}]

        resp = {
            "ok": True,
            "hash": sha,
            "meta": {
                "pages": 1,
                "rows": len(rows),
                "columns": len(headers) if headers else 0,
                "lang": "auto",
                "cache_hit": False,
            },
            "confidence": 1.0,
            "pages": pages,
            "full_text": full_text,
            "warnings": [],
        }
        write_cache(sha, resp)
        return resp
    except Exception:
        return error_response(400, "E_CSV_OPEN", "No se pudo procesar el archivo CSV")


@app.post("/extract/txt")
def extract_txt(file: UploadFile = File(...)) -> Dict[str, Any]:
    if file.content_type not in ALLOWED_TXT_TYPES:
        return error_response(
            415, "E_UNSUPPORTED_TYPE", f"Tipo no soportado: {file.content_type}"
        )
    content = read_upload(file)
    sha = compute_sha256(content)
    cached = load_cache(sha)
    if cached:
        cached.setdefault("meta", {})["cache_hit"] = True
        return cached
    try:
        # Detectar encoding
        detected = chardet.detect(content)
        encoding = detected.get("encoding", "utf-8")

        full_text = content.decode(encoding)
        lines = full_text.split("\n")

        pages = [{"number": 1, "text": full_text, "confidence": 1.0}]

        resp = {
            "ok": True,
            "hash": sha,
            "meta": {
                "pages": 1,
                "lines": len(lines),
                "encoding": encoding,
                "lang": "auto",
                "cache_hit": False,
            },
            "confidence": 1.0,
            "pages": pages,
            "full_text": full_text,
            "warnings": [],
        }
        write_cache(sha, resp)
        return resp
    except Exception:
        return error_response(
            400, "E_TXT_OPEN", "No se pudo procesar el archivo de texto"
        )


@app.post("/extract/pptx")
def extract_pptx(file: UploadFile = File(...)) -> Dict[str, Any]:
    if file.content_type not in ALLOWED_PPTX_TYPES:
        return error_response(
            415, "E_UNSUPPORTED_TYPE", f"Tipo no soportado: {file.content_type}"
        )
    content = read_upload(file)
    sha = compute_sha256(content)
    cached = load_cache(sha)
    if cached:
        cached.setdefault("meta", {})["cache_hit"] = True
        return cached
    try:
        pptx_file = io.BytesIO(content)
        presentation = Presentation(pptx_file)

        slides_text = []
        full_text_parts = []

        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = []

            # Extraer texto de todas las formas en la diapositiva
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

                # Si es una tabla, extraer el texto de las celdas
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_text.append(" | ".join(row_text))

            slide_content = "\n".join(slide_text)
            slides_text.append(
                {"number": slide_num, "text": slide_content, "confidence": 1.0}
            )

            if slide_content.strip():
                full_text_parts.append(
                    f"=== DIAPOSITIVA {slide_num} ===\n{slide_content}"
                )

        full_text = "\n\n".join(full_text_parts)

        resp = {
            "ok": True,
            "hash": sha,
            "meta": {
                "slides": len(presentation.slides),
                "pages": len(presentation.slides),
                "lang": "auto",
                "cache_hit": False,
            },
            "confidence": 1.0,
            "pages": slides_text,
            "full_text": full_text,
            "warnings": [],
        }
        write_cache(sha, resp)
        return resp
    except Exception as e:
        return error_response(
            400, "E_PPTX_OPEN", f"No se pudo procesar el archivo PowerPoint: {str(e)}"
        )


# Hooks utilitarios para futuro fallback a la nube (no implementados)
def send_to_cloud_fallback(
    _sha: str, _payload: Dict[str, Any]
) -> None:  # pragma: no cover
    # Placeholder para integrar con un servicio remoto si se desea
    return None


# =============================
# OCR especializados (modulares)
# =============================


@app.post("/ocr/nums")
def ocr_only_numbers(file: UploadFile = File(...)) -> Dict[str, Any]:
    if file.content_type not in ALLOWED_IMAGE_TYPES:
        return error_response(
            415, "E_UNSUPPORTED_TYPE", f"Tipo no soportado: {file.content_type}"
        )
    content = read_upload(file)
    try:
        from PIL import Image

        img = Image.open(io.BytesIO(content))
        # Temporary: use normal OCR for numbers
        ocr_result = ocr_image_to_text(img)
        text_content = ocr_result if isinstance(ocr_result, str) else ocr_result.get("text", "")
        res = {"text": text_content, "numbers": []}
        ft = (res.get("text") or " ".join(res.get("numbers", []) or []) or "").strip()
        return {"ok": True, "full_text": ft, **res}
    except Exception as e:
        return error_response(400, "E_NUMS", f"No se pudo reconocer números: {str(e)}")


@app.post("/ocr/math")
def ocr_math(file: UploadFile = File(...)) -> Dict[str, Any]:
    if file.content_type not in ALLOWED_IMAGE_TYPES:
        return error_response(
            415, "E_UNSUPPORTED_TYPE", f"Tipo no soportado: {file.content_type}"
        )
    content = read_upload(file)
    try:
        from PIL import Image

        img = Image.open(io.BytesIO(content))
        # Temporary: use normal OCR for math
        ocr_result = ocr_image_to_text(img)
        text_content = ocr_result if isinstance(ocr_result, str) else ocr_result.get("text", "")
        res = {"latex": text_content}
        ft = (res.get("latex") or "").strip()
        return {"ok": True, "full_text": ft, **res}
    except Exception as e:
        return error_response(400, "E_MATH", f"No se pudo reconocer fórmula: {str(e)}")


@app.post("/ocr/mixed")
def ocr_mixed(file: UploadFile = File(...)) -> Dict[str, Any]:
    if file.content_type not in ALLOWED_IMAGE_TYPES:
        return error_response(
            415, "E_UNSUPPORTED_TYPE", f"Tipo no soportado: {file.content_type}"
        )

    try:
        content = read_upload(file)
        print(
            f"[OCR] File received: {file.filename}, size: {len(content)} bytes, type: {file.content_type}"
        )
    except Exception as e:
        print(f"[OCR ERROR] Failed to read upload: {e}")
        return error_response(400, "E_READ_UPLOAD", f"Error leyendo archivo: {str(e)}")

    try:
        from PIL import Image

        img = Image.open(io.BytesIO(content))
        print(f"[OCR] Image opened successfully: {img.size} pixels, mode: {img.mode}")

        ocr_result = ocr_image_to_text(img)
        text_content = ocr_result if isinstance(ocr_result, str) else ocr_result.get("text", "")
        res = {
            "markdown": text_content,
            "text": text_content,
            "provider": "ensemble",
        }
        ft = (res.get("markdown") or res.get("text") or "").strip()
        provider = res.get("provider")

        # Log SIEMPRE visible
        try:
            preview = ft[:80] + ("…" if len(ft) > 80 else "")
            print(
                f"[OCR] /ocr/mixed provider={provider} conf={res.get('confidence','-')} text='{preview}'"
            )
        except Exception:
            pass

        return {"ok": True, "full_text": ft, **res}

    except Exception as e:
        print(f"[OCR ERROR] Processing failed: {e}")
        import traceback

        traceback.print_exc()
        return error_response(
            400, "E_MIXED", f"No se pudo reconocer contenido mixto: {str(e)}"
        )


@app.post("/ocr/chem")
def ocr_chem(file: UploadFile = File(...)) -> Dict[str, Any]:
    if file.content_type not in ALLOWED_IMAGE_TYPES:
        return error_response(
            415, "E_UNSUPPORTED_TYPE", f"Tipo no soportado: {file.content_type}"
        )
    content = read_upload(file)
    try:
        from PIL import Image

        img = Image.open(io.BytesIO(content))
        # Temporary: use normal OCR for chemistry
        ocr_result = ocr_image_to_text(img)
        text_content = ocr_result if isinstance(ocr_result, str) else ocr_result.get("text", "")
        res = {"smiles": text_content}
        ft = (res.get("smiles") or "").strip()
        return {"ok": True, "full_text": ft, **res}
    except Exception as e:
        return error_response(
            400, "E_CHEM", f"No se pudo reconocer estructura química: {str(e)}"
        )


@app.post("/ocr/phys_post")
def ocr_phys_post(file: UploadFile = File(...)) -> Dict[str, Any]:
    """Endpoint auxiliar: extrae texto y post-procesa unidades/expresiones físicas."""
    if file.content_type not in ALLOWED_IMAGE_TYPES:
        return error_response(
            415, "E_UNSUPPORTED_TYPE", f"Tipo no soportado: {file.content_type}"
        )
    content = read_upload(file)
    try:
        from PIL import Image

        img = Image.open(io.BytesIO(content))
        base = ocr_image_to_text(img)
        # Temporary: skip physics postprocessing
        base_text = base if isinstance(base, str) else base.get("text", "")
        pp = base_text
        return {
            "ok": True,
            "text": base_text,
            "full_text": base_text,
            "physics": pp,
        }
    except Exception as e:
        return error_response(
            400, "E_PHYS", f"No se pudo post-procesar física: {str(e)}"
        )
